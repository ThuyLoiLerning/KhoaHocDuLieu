"""Khảo sát từng slide PPTX: đo textbox có tràn khung không, mật độ chữ, bố cục.

Ước lượng số dòng wrap: ~1.9 ký tự/inch/pt → chars_per_line = width_in * font_pt * 1.9 / 17? Dùng xấp xỉ đơn giản:
chars_per_line ≈ width_inches * 2.0 * (17 / font_pt). Chiều cao dòng ≈ font_pt * 1.25 / 72 inch.
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu

PPTX = Path("reports/slides/TrinhBay_ChuyenDe4.pptx")
prs = Presentation(PPTX)

def emu_in(v):
    return v / 914400.0

issues = []
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        text = tf.text.strip()
        if not text:
            continue
        w = emu_in(sh.width)
        h = emu_in(sh.height)
        top = emu_in(sh.top)
        left = emu_in(sh.left)
        # Ước lượng chiều cao needed cho mỗi paragraph
        needed = 0.0
        n_lines_est = 0
        for p in tf.paragraphs:
            pt = 17.0 if p.level == 0 else 15.0
            for r in p.runs:
                if r.font.size:
                    pt = r.font.size.pt
            t = p.text
            cpl = max(8, w * 144.0 / pt)  # ký tự/dòng xấp xỉ: w_in*96px/(0.5em×pt×1.333)
            lines = max(1, int(len(t) / cpl) + (1 if len(t) % cpl else 0))
            n_lines_est += lines
            needed += lines * pt * 1.25 / 72.0 + (0.11 if p.level == 0 else 0.055)
        bottom = top + h
        avail_ok = needed <= h + 0.15
        flag = "" if avail_ok else " <<< CÓ THỂ TRÀN"
        if not avail_ok:
            issues.append(f"  slide {i}: textbox '{text[:35]}...' cần {needed:.2f}in > khung {h:.2f}in{flag}")
        texts.append(f"{text[:28]}(w={w:.1f},h={h:.1f})")
    # Số key per slide
    n_txt = len(texts)
    n_tbl = sum(1 for sh in slide.shapes if sh.has_table)
    n_pic = sum(1 for sh in slide.shapes if sh.shape_type == 13)
    print(f"--- Slide {i}: textbox={n_txt} table={n_tbl} pic={n_pic}")
    for t in texts:
        print(f"    {t}")

if issues:
    print("\n=== NGHI NGỜ TRÀN ===")
    for it in issues:
        print(it)
else:
    print("\nKhông phát hiện tràn theo ước lượng.")