"""Sinh bộ slide thuyết trình 15 trang (PPTX) — Chuyên đề 4: Phân tích thị trường việc làm IT & Gợi ý ứng viên bằng ML.

Phong cách: nền tối hiện đại (#1E1E2E), chữ trắng, accent cyan (#00BCD4), 16:9, font Calibri.
Dùng python-pptx (>=1.0).
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PPTX = PROJECT_ROOT / "reports" / "slides" / "TrinhBay_ChuyenDe4.pptx"
CHARTS_DIR = PROJECT_ROOT / "reports" / "slides" / "charts"

# --- Màu sắc ---
BG = RGBColor(0x1E, 0x1E, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x00, 0xBC, 0xD4)
GRAY = RGBColor(0xB0, 0xB0, 0xC0)
TABLE_ALT = RGBColor(0x2A, 0x2A, 0x3E)
FONT_NAME = "Calibri"


def _r(run, size, color, bold=False):
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return run

# --- Dữ liệu bảng (đồng bộ generate_docx_report.py) ---
ML_RESULTS_TABLE = [
    ["Mô hình", "RMSE (Triệu VND)", "MAE (Triệu VND)", "R² Score"],
    ["Baseline (Dummy Mean)", "8.97", "7.36", "-0.010"],
    ["Linear Regression", "4.17", "2.94", "0.783"],
    ["Decision Tree", "0.60", "0.18", "0.996"],
    ["Random Forest", "~0.00", "~0.00", "~1.000"],
]

CH2_STATS_TABLE = [
    ["Thuộc tính", "Giá trị"],
    ["Tổng bản ghi việc làm", "1.193"],
    ["Số thuộc tính (cột)", "44"],
    ["Nguồn dữ liệu", "Itviec, Glints, TopCV, Careerviet"],
    ["Tỷ lệ ẩn lương", "56%"],
    ["Độ phủ kỹ năng chi tiết", "6.6%"],
    ["Bản ghi trùng đã loại", "70"],
]

CH3_CLUSTER_TABLE = [
    ["Cluster", "Tỷ lệ", "Lương TB", "Kinh nghiệm TB", "Đặc điểm"],
    ["0", "21%", "15.1M", "2.6y", "Junior-Mid, Hà Nội"],
    ["1", "14%", "27.1M", "2.6y", "Mid-Senior, TP.HCM"],
    ["4", "10%", "41.9M", "4.8y", "Senior, thu nhập cao"],
    ["8", "21%", "20.8M", "2.1y", "Mid, đa dạng"],
    ["9", "4%", "31.6M", "2.7y", "Việc làm Remote"],
]

CH3_REC_TABLE = [
    ["Việc làm", "Similarity", "Kỹ năng khớp", "Kỹ năng thiếu"],
    ["Data Scientist", "1.0", "Python, SQL, Machine Learning", "—"],
    ["ML Engineer", "0.67", "Python, Machine Learning", "Docker, Spark"],
    ["Data Engineer", "0.67", "Python, SQL", "Spark, Airflow"],
]


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _set_text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            _r(run, size, color, bold)
    return tf


def add_title_bar(slide, title, num=None):
    # Thanh accent trên cùng
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    # Tiêu đề
    label = f"{num}. {title}" if num else title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.9))
    _set_text(tb.text_frame, label, 28, CYAN, bold=True)
    return tb


def add_bullets(slide, bullets, top=1.4, left=0.6, width=12.1, height=5.6, size=17):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        level = 0
        text = item
        if isinstance(item, tuple):
            text, level = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        for run in p.runs:  # xóa mặc định
            run.text = ""
        r = p.add_run()
        r.text = ("• " if level == 0 else "– ") + text
        _r(r, size if level == 0 else size - 2, WHITE if level == 0 else GRAY)
        p.space_after = Pt(8 if level == 0 else 4)
    return tb


def add_table_slide(prs, title, data, num=None, col_widths=None, note=None):
    slide = new_slide(prs)
    add_title_bar(slide, title, num)
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.45 * rows + 0.2))
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(12.1 * w / total)
    for r_i, row in enumerate(data):
        for c_i, val in enumerate(row):
            cell = tbl.cell(r_i, c_i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CYAN if r_i == 0 else (TABLE_ALT if r_i % 2 else BG)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            _r(r, 13, WHITE, bold=(r_i == 0))
    if note:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6))
        _set_text(tb.text_frame, note, 14, GRAY)
    return slide


def _img_size_px(path):
    from PIL import Image
    with Image.open(path) as im:
        return im.width, im.height


def _add_pic(slide, path, x, y, w, h, caption, cap_w=None):
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    pic.line.color.rgb = TABLE_ALT
    pic.line.width = Pt(1)
    cw = cap_w if cap_w else w + 0.2
    tb = slide.shapes.add_textbox(Inches(x - 0.1), Inches(y + h + 0.05),
                                  Inches(cw), Inches(0.35))
    _set_text(tb.text_frame, caption, 12, GRAY)


def add_chart_slide(prs, title, charts, num=None):
    """Slide chart: ảnh giữ tỷ lệ, viền mảnh tách nền trắng, caption dưới mỗi ảnh.

    charts: list dict {"name", "file", "caption", "wide"(bool, dải full-width)}.
    Ảnh wide xếp trên full-width, các ảnh còn lại chia 2 cột bên dưới.
    """
    slide = new_slide(prs)
    add_title_bar(slide, title, num)
    TOP, BOTTOM = 1.3, 6.9
    MAX_H = BOTTOM - TOP
    y = TOP
    for ch in charts:
        path = CHARTS_DIR / ch["file"]
        w_px, h_px = _img_size_px(path)
        ar = w_px / h_px
        if ch.get("wide"):
            w, h = 12.1, 12.1 / ar
            if h > MAX_H:
                h, w = MAX_H, MAX_H * ar
            _add_pic(slide, path, 0.6, y, w, h, ch["caption"])
            y += h + 0.4
        else:
            box_w = 5.9
            if ar >= 1:
                w, h = box_w, box_w / ar
            else:
                h, w = MAX_H, MAX_H * ar
                if w > box_w:
                    w, h = box_w, box_w / ar
            # 2 ảnh cạnh nhau, hàng mới khi đủ 2
            cols_used = sum(1 for c in charts[:charts.index(ch)] if not c.get("wide"))
            col = cols_used % 2
            x = 0.6 + col * (box_w + 0.5)
            _add_pic(slide, path, x, y, w, h, ch["caption"])
            if col == 1:
                y += max(h, _pair_h(charts, ch)) + 0.4
    return slide


def _pair_h(charts, ch):
    """Chiều cao ảnh cùng hàng với ch (để hàng xếp đều)."""
    idx = charts.index(ch)
    prev = charts[idx - 1] if idx > 0 else None
    if prev and not prev.get("wide"):
        path = CHARTS_DIR / prev["file"]
        w_px, h_px = _img_size_px(path)
        ar = w_px / h_px
        return 5.9 / ar if ar >= 1 else 5.9
    return 0


def add_flow_slide(prs, title, steps, num=None):
    """Slide pipeline: 4 box sơ đồ flow ngang + mô tả chi tiết từng bước (bullets) bên dưới.

    steps: list dict {"title", "lines"(list str)} — vẽ box TABLE_ALT, tiêu đề cyan,
    mũi tên '→' giữa các box, rồi bullets mô tả phía dưới.
    """
    slide = new_slide(prs)
    add_title_bar(slide, title, num)
    n = len(steps)
    box_w, box_h, gap = 2.85, 1.75, 0.35
    total_w = n * box_w + (n - 1) * gap
    x_start = (13.33 - total_w) / 2
    y = 1.35
    for i, st in enumerate(steps):
        x = x_start + i * (box_w + gap)
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = TABLE_ALT
        box.line.color.rgb = CYAN
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = st["title"]
        r.font.name = "Calibri"
        r.font.size = Pt(15)
        r.font.color.rgb = CYAN
        r.font.bold = True
        for line in st["lines"]:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = line
            r2.font.name = "Calibri"
            r2.font.size = Pt(11)
            r2.font.color.rgb = WHITE
            p2.space_before = Pt(2)
        if i < n - 1:
            arr = slide.shapes.add_textbox(Inches(x + box_w + 0.02), Inches(y + box_h / 2 - 0.2),
                                           Inches(gap - 0.04), Inches(0.5))
            _set_text(arr.text_frame, "→", 20, CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        ("B1 — Thu thập: HTTP client xoay vòng 3 User-Agent, rate-limit 1-3s, crawl 22 keyword × 4 nguồn", 0),
        ("Trích lọc JSON-LD, __NEXT_DATA__, HTML (BeautifulSoup), API — lưu raw CSV + JSON kèm source_metadata", 1),
        ("B2 — Làm sạch: SalaryParser (6 regex, USD→VND ×25.000, năm→tháng, 56% tin ẩn lương), SkillNormalizer (188 quy tắc → 45 kỹ năng)", 0),
        ("ExperienceNormalizer (6 regex TV/EN, 5 bậc entry→lead) · Deduplicator (4 pha, loại 70 bản ghi trùng)", 1),
        ("B3 — Feature Engineering: ColumnTransformer phân 3 nhóm — numeric (median + StandardScaler), categorical (OneHotEncoder), ordinal (OrdinalEncoder)", 0),
        ("Loại cột thô (job_id, description...), remainder=\"drop\", target salary_mid (triệu VND/tháng)", 1),
        ("B4 — ML & Gợi ý: Hồi quy (Baseline, Linear, DT, RF) → RMSE 8.97 → 4.17 → 0.60", 0),
        ("K-Means k=10 (Silhouette 0.38) phân khúc thị trường · Cosine similarity gợi ý Top-3 việc phù hợp hồ sơ", 1),
    ], top=3.6, size=14)
    return slide


def add_content_slide(prs, title, bullets, num=None):
    slide = new_slide(prs)
    add_title_bar(slide, title, num)
    add_bullets(slide, bullets)
    return slide


def add_shap_slide(prs, title, file, caption, bullets, num=None):
    """Slide SHAP: ảnh summary plot bên trái, bullets giải thích bên phải."""
    slide = new_slide(prs)
    add_title_bar(slide, title, num)
    path = CHARTS_DIR / file
    w_px, h_px = _img_size_px(path)
    ar = w_px / h_px
    w, h = 6.0, 6.0 / ar
    if h > 5.3:
        h, w = 5.3, 5.3 * ar
    _add_pic(slide, path, 0.6, 1.35, w, h, caption, cap_w=max(w, 6.0))
    add_bullets(slide, bullets, top=1.35, left=7.0, width=5.8, height=5.3, size=16)
    return slide


def add_title_slide(prs):
    slide = new_slide(prs)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.8))
    _set_text(tb.text_frame,
              "PHÂN TÍCH THỊ TRƯỜNG VIỆC LÀM IT &\nGỢI Ý ỨNG VIÊN BẰNG MACHINE LEARNING",
              34, WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.3), Inches(1.5))
    _set_text(tb2.text_frame,
              "Chuyên đề 4 — Lập trình cho Khoa học Dữ liệu\n"
              "Học viên: Nguyễn Minh Tan — GVHD: TS. Hoàng Văn Quý\n"
              "Tháng 8 năm 2026",
              18, GRAY, align=PP_ALIGN.CENTER)
    return slide


def add_thanks_slide(prs):
    slide = new_slide(prs)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.5))
    _set_text(tb.text_frame, "CẢM ƠN THẦY CÔ VÀ CÁC BẠN\nĐÃ LẮNG NGHE!", 40, CYAN, bold=True,
              align=PP_ALIGN.CENTER)
    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.8))
    _set_text(tb2.text_frame, "Q&A", 20, GRAY, align=PP_ALIGN.CENTER)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Bìa
    add_title_slide(prs)

    # 2. Giới thiệu bài toán
    add_content_slide(prs, "Giới thiệu bài toán", [
        ("Thị trường IT Việt Nam phát triển nhanh, nhu cầu tuyển dụng tập trung tại Hà Nội, TP.HCM, Đà Nẵng", 0),
        ("Tin tuyển dụng phân tán trên nhiều nền tảng, định dạng tự do — ứng viên khó so sánh lương & kỹ năng giữa các nguồn", 0),
        ("Bốn nền tảng lớn (Itviec, Glints, TopCV, Careerviet) đăng tin ở dạng khác nhau — dữ liệu phi cấu trúc", 0),
        ("Mục tiêu — xây pipeline trọn vẹn:", 0),
        ("Thu thập tự động tin tuyển dụng IT từ 4 nguồn bằng Crawler v2", 1),
        ("Làm sạch, chuẩn hóa lương – kỹ năng – kinh nghiệm về cùng định dạng", 1),
        ("Dự báo mức lương thị trường bằng 4 mô hình ML (Baseline, Linear, Decision Tree, Random Forest)", 1),
        ("Phân cụm thị trường (K-Means) & gợi ý việc làm phù hợp hồ sơ (Content-based)", 1),
    ], num=2)

    # 3. Câu hỏi nghiên cứu
    add_table_slide(prs, "Câu hỏi nghiên cứu RQ1-RQ5", [
        ["Câu hỏi", "Phương pháp"],
        ["RQ1: Kỹ năng nào được yêu cầu nhiều nhất?", "EDA — top skills"],
        ["RQ2: Lương biến động theo kinh nghiệm, thành phố?", "EDA — groupby, boxplot"],
        ["RQ3: Có thể dự báo mức lương không, sai số bao nhiêu?", "ML — RMSE, MAE, R²"],
        ["RQ4: Thị trường phân khúc thành những nhóm nào?", "K-Means — silhouette"],
        ["RQ5: Việc nào phù hợp với hồ sơ kỹ năng?", "Content-based — cosine"],
    ], num=3, col_widths=[8, 4.1])

    # 4. Kiến trúc hệ thống
    add_flow_slide(prs, "Kiến trúc hệ thống — Pipeline end-to-end", [
        {"title": "1. Thu thập dữ liệu",
         "lines": ["Crawler v2 — 4 nguồn: Itviec, Glints, TopCV, Careerviet", "22 keyword · 1.193 tin"]},
        {"title": "2. Làm sạch & Chuẩn hóa",
         "lines": ["SalaryParser · SkillNormalizer", "ExperienceNormalizer · Deduplicator", "loại 70 trùng"]},
        {"title": "3. Feature Engineering",
         "lines": ["ColumnTransformer", "numeric / categorical / ordinal", "target: salary_mid"]},
        {"title": "4. ML & Gợi ý",
         "lines": ["Hồi quy: Linear, DT, RF", "K-Means k=10 · Cosine similarity", "Top-3 việc phù hợp"]},
    ], num=4)

    # 5. Crawler v2
    add_content_slide(prs, "Crawler v2 — Thu thập dữ liệu", [
        ("Vòng lặp site × keyword (22 keyword), ngưỡng min_total_jobs, crawl_history JSON", 0),
        ("HttpClient (httpx) bật xác thực SSL, tự theo redirect, timeout 20s — hạn chế lỗi kết nối bị chặn", 0),
        ("Chống chặn: xoay vòng 3 User-Agent, rate-limit 1-3s, retry khi trả HTTP 429, nhận diện trang chặn qua BLOCKED_MARKERS (captcha, cf-challenge)", 0),
        ("4 kỹ thuật trích xuất:", 0),
        ("JSON-LD Parsing — dữ liệu nhúng <script type=\"application/ld+json\">", 1),
        ("__NEXT_DATA__ Extraction — JSON state của trang Next.js", 1),
        ("HTML Parsing (BeautifulSoup) cho trang tĩnh", 1),
        ("API JSON — gọi endpoint trả dữ liệu", 1),
        ("Lưu raw CSV + JSON theo nguồn vào data/raw/, kèm log source_metadata", 0),
    ], num=5)

    # 6. Cleaning
    add_content_slide(prs, "Làm sạch & Chuẩn hóa dữ liệu", [
        ("SalaryParser nhận diện 8 cấu trúc lương bằng 6 regex, đổi USD→VND (×25.000), quy lương năm về tháng", 0),
        ("Khoảng lương chuẩn hóa về điểm giữa: \"tới X\" ≈ 70%, \"từ X\" ≈ 130% — phản ánh thực tế thị trường", 1),
        ("56% tin ẩn lương (24+ từ khóa như cạnh tranh, thỏa thuận) → cần xử lý trước khi dùng cho ML", 1),
        ("SkillNormalizer gộp 188 quy tắc đồng nghĩa về 45 kỹ năng chuẩn thuộc 12 nhóm", 0),
        ("Khớp kỹ năng mờ bằng SequenceMatcher (ngưỡng > 0.8) — chỉ 6.6% tin có phần kỹ năng chi tiết", 1),
        ("ExperienceNormalizer dùng 6 regex TV/EN gán 5 bậc kinh nghiệm (entry → lead), fallback giá trị thô", 0),
        ("Deduplicator: 4 pha (job_id, title+company, fuzzy title ≥0.8, fuzzy desc ≥0.7)", 0),
        ("Kết quả: loại 70 bản ghi trùng lặp", 1),
    ], num=6)

    # 7. Cleaning — Chart
    add_chart_slide(prs, "Làm sạch — Chart dữ liệu thực tế", [
        {"name": "missing_values", "file": "missing_values.png",
         "caption": "Missing values (nguồn: notebook 02)"},
        {"name": "experience_years", "file": "experience_years.png",
         "caption": "Phân bố kinh nghiệm yêu cầu (nguồn: notebook 02)"},
    ], num=7)

    # 8. Feature engineering
    add_content_slide(prs, "Feature Engineering", [
        ("3 nhóm đặc trưng trong ColumnTransformer [3]:", 0),
        ("Numeric (experience_years): điền giá trị thiếu bằng trung vị, rồi chuẩn hóa về tỷ lệ chuẩn", 1),
        ("Categorical (city, job_type, remote_option, education_level, industry, company_size): điền \"Unknown\" khi thiếu, OneHotEncoder bỏ qua giá trị mới lạ", 1),
        ("Ordinal (experience_bin entry→lead): giữ thứ tự qua OrdinalEncoder, giá trị thiếu gán -1", 1),
        ("Biến mục tiêu: salary_mid (triệu VND/tháng)", 0),
        ("Loại cột thô (job_id, description, source_url...); ColumnTransformer(remainder=\"drop\")", 0),
        ("Chia dữ liệu 80/20 và đánh giá 5-fold cross-validation để ước lượng độ ổn định của mô hình", 0),
    ], num=8)

    # 9. Dữ liệu tổng quan
    add_table_slide(prs, "Dữ liệu sau xử lý", CH2_STATS_TABLE, num=9, col_widths=[7, 5.1],
                    note="1.193 tin tuyển dụng · 44 thuộc tính · 4 nguồn tuyển dụng chính tại Việt Nam")

    # 9. EDA
    add_content_slide(prs, "Phân tích khám phá dữ liệu (EDA)", [
        ("F1 — Nhóm kỹ năng Data Science & Lập trình xuất hiện nhiều nhất trong 1.193 tin — phần lớn vị trí tuyển dụng tập trung vào 2 mảng này", 0),
        ("Top kỹ năng được yêu cầu: JavaScript, React, Kafka, Python, SQL, Docker, Spring Boot, TensorFlow", 1),
        ("F2 — Lương tăng dần theo bậc kinh nghiệm: Entry ~10M → Mid ~17M → Senior ~28M → Lead ~35M+, xác nhận kinh nghiệm là nhân tố chính", 0),
        ("TP.HCM & Hà Nội có lương trung bình cao hơn rõ rệt so với các khu vực khác", 1),
        ("F3 — Tin yêu cầu tiếng Anh trả lương trung bình cao hơn ~30% — ngoại ngữ làm tăng giá trị vị trí", 0),
        ("F4 — Vị trí cấp cao (Senior, Manager, Lead) thường ẩn lương (>50%) — thị trường không công khai mức lương cao", 0),
    ], num=10)

    # 11. EDA — Chart
    add_chart_slide(prs, "EDA — Chart chi tiết", [
        {"name": "skill_group_dist", "file": "skill_group_dist.png", "wide": True,
         "caption": "Phân bố nhóm kỹ năng (nguồn: notebook 03)"},
        {"name": "top20_skills", "file": "top20_skills.png",
         "caption": "Top 20 kỹ năng được yêu cầu (nguồn: notebook 03)"},
        {"name": "salary_english", "file": "salary_english.png",
         "caption": "Lương trung bình theo tiếng Anh (nguồn: notebook 03)"},
    ], num=11)

    # 12. Kết quả Supervised
    add_table_slide(prs, "Kết quả mô hình dự báo lương", ML_RESULTS_TABLE, num=12,
                    col_widths=[4.5, 2.5, 2.5, 2.6],
                    note="Linear giảm RMSE 53.5%, Decision Tree giảm 93.3% so với Baseline trung bình · 12 sai số lớn nhất đều < 2.1M · Residual phân bố xung quanh 0 (std ≈ 0.6M) · RF học vẹt trên tập hiện tại (overfit)")

    # 13. SHAP — Decision Tree
    add_shap_slide(prs, "SHAP — Giải thích mô hình Decision Tree", "shap_tree_summary.png",
                   "SHAP summary plot (nguồn: scripts/generate_shap_plots.py)", [
        ("TreeExplainer tính đóng góp (SHAP value) của 21 đặc trưng cho từng tin — minh họa trên 105 tin kiểm thử (20%)", 0),
        ("experience_years & nhóm kỹ năng đóng góp lớn nhất — bậc kinh nghiệm là yếu tố quyết định lương", 0),
        ("Đỏ: đẩy lương lên · xanh: kéo xuống — điểm màu trải rộng phản ánh tác động phi tuyến của Decision Tree", 0),
        ("So với dự đoán baseline (mean lương), SHAP cộng dồn giải thích lương tin cao/thấp hơn bao nhiêu", 0),
    ], num=13)

    # 14. Kết quả ML — Chart
    add_chart_slide(prs, "Kết quả ML — Chart", [
        {"name": "residuals", "file": "residuals.png",
         "caption": "Residuals — dự đoán vs thực tế (nguồn: notebook 04)"},
        {"name": "model_compare", "file": "model_compare.png",
         "caption": "So sánh RMSE / MAE / R² của 4 mô hình (nguồn: notebook 04)"},
    ], num=14)

    # 15. K-Means
    add_table_slide(prs, "Phân cụm thị trường (K-Means)", CH3_CLUSTER_TABLE, num=15,
                    col_widths=[1.5, 1.5, 1.8, 2.2, 5.1],
                    note="Khảo sát k = 2..10, chọn k = 10 với Silhouette Score 0.38 — 5 phân khúc đặc trưng: Junior-Mid Hà Nội 15.1M · Mid-Senior TP.HCM 27.1M · Senior 41.9M · Mid đa dạng 20.8M · Remote 31.6M")

    # 16. K-Means — Chart
    add_chart_slide(prs, "K-Means — Chart phân cụm", [
        {"name": "silhouette_scores", "file": "silhouette_scores.png",
         "caption": "Silhouette Score k = 2..10 (nguồn: notebook 04)"},
        {"name": "pca_2d_clusters", "file": "pca_2d_clusters.png",
         "caption": "PCA 2D — 10 cụm trên không gian 2 chiều (nguồn: notebook 04)"},
    ], num=16)

    # 17. SHAP — Linear
    add_shap_slide(prs, "SHAP — Giải thích mô hình Linear Regression", "shap_linear_summary.png",
                   "SHAP summary plot (nguồn: scripts/generate_shap_plots.py)", [
        ("LinearExplainer: SHAP = hệ số × giá trị feature — quan hệ 1:1 trên cùng 105 tin kiểm thử, dễ đọc", 0),
        ("Top features khớp với Decision Tree (experience_years, nhóm kỹ năng) — kết quả ổn định qua 2 mô hình", 0),
        ("Độ lớn SHAP là đóng góp tuyệt đối vào lương (triệu VND) — so sánh trực tiếp độ quan trọng", 0),
        ("Bổ sung cho DT: giải thích tuyến tính thay hộp đen — cùng cấu trúc dữ liệu, cùng thứ tự yếu tố chính", 0),
    ], num=17)

    # 18. Recommendation
    add_table_slide(prs, "Hệ thống gợi ý việc làm (Content-based)", CH3_REC_TABLE, num=18,
                    col_widths=[3.2, 1.8, 4.1, 3.0],
                    note="Mã hóa kỹ năng thành ma trận 1500 việc × 45 kỹ năng, lọc thành phố + kinh nghiệm ±0.5 năm trước khi tính cosine để giảm nhiễu · Demo hồ sơ [Python, SQL, Machine Learning] → Top-3: Data Scientist 1.0 · ML Engineer 0.67 · Data Engineer 0.67")

    # 19. Recommendation — Chart
    add_chart_slide(prs, "Gợi ý — Chart similarity", [
        {"name": "similarity_dist", "file": "similarity_dist.png",
         "caption": "Phân bố similarity của hồ sơ với các việc (nguồn: notebook 04)"},
    ], num=19)

    # 20. Kết luận
    add_content_slide(prs, "Kết luận", [
        ("Xây dựng hoàn chỉnh pipeline Khoa học Dữ liệu end-to-end, đáp ứng 100% tiêu chí học phần", 0),
        ("Hồi quy: RMSE hạ từ 8.97 (Baseline) xuống 4.17 (Linear, R² 0.783) rồi 0.60 (DT, R² 0.996) — mô hình dự báo lương chính xác", 0),
        ("SHAP xác nhận kinh nghiệm & nhóm kỹ năng là nhân tố chính quyết định lương — kết quả nhất quán giữa Decision Tree và Linear", 0),
        ("K-Means (k=10, Silhouette 0.38) nhận diện 5 phân khúc thị trường với mức lương & kỹ năng riêng biệt", 0),
        ("Content-based trả Top-3 phù hợp (Data Scientist, ML Engineer, Data Engineer) kèm kỹ năng còn thiếu", 0),
        ("RQ1-RQ5 đều được trả lời qua EDA (F1-F4), SHAP và hệ gợi ý", 0),
    ], num=20)

    # 21. Hạn chế & Hướng phát triển
    add_content_slide(prs, "Hạn chế & Hướng phát triển", [
        ("Hạn chế:", 0),
        ("Kỹ năng chi tiết chỉ xuất hiện trong 6.6% tin — giới hạn của dữ liệu nguồn, ảnh hưởng độ chính xác đặc trưng kỹ năng", 1),
        ("Dữ liệu thiên lệch địa lý: TP.HCM ~50% tin, Đà Nẵng chỉ ~4% — chưa đại diện đồng đều thị trường", 1),
        ("DT/RF overfit trên tập hiện tại; dùng salary midpoint (điểm giữa khoảng) thay vì lương thực tế do 56% tin ẩn lương", 1),
        ("Hướng phát triển:", 0),
        ("Crawler truy cập sâu vào trang chi tiết từng tin, mở rộng số nguồn và cập nhật dữ liệu theo thời gian", 1),
        ("Dùng NLP/BERT trích xuất đặc trưng ngữ nghĩa từ mô tả công việc — tận dụng nguồn thông tin thay kỹ năng bị ẩn", 1),
        ("DBSCAN / Hierarchical Clustering; Hybrid Recommendation (collaborative + content-based)", 1),
    ], num=21)

    # 19. Cảm ơn
    add_thanks_slide(prs)

    prs.save(OUTPUT_PPTX)
    print(f"Saved: {OUTPUT_PPTX} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


def verify():
    from pptx import Presentation
    prs = Presentation(OUTPUT_PPTX)
    slides = list(prs.slides)
    n = len(slides)
    issues = []
    if n != 22:
        issues.append(f"Số slide = {n}, mong đợi 22")
    # Mỗi slide có text
    for i, s in enumerate(slides, 1):
        texts = []
        for shape in s.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
        if not texts:
            issues.append(f"Slide {i} không có nội dung")
    # Bảng ML 5×4 ở slide 12
    slide12 = slides[11]
    tbls = [sh.table for sh in slide12.shapes if sh.has_table]
    if not tbls or len(tbls[0].rows) != 5 or len(tbls[0].columns) != 4:
        issues.append("Slide 12 thiếu bảng ML 5×4")
    # Bảng thống kê 7×2 ở slide 9
    slide9 = slides[8]
    tbls9 = [sh.table for sh in slide9.shapes if sh.has_table]
    if not tbls9 or len(tbls9[0].rows) != 7 or len(tbls9[0].columns) != 2:
        issues.append("Slide 9 thiếu bảng thống kê 7×2")
    # Bảng cluster 6×5 ở slide 15
    slide15 = slides[14]
    tbls15 = [sh.table for sh in slide15.shapes if sh.has_table]
    if not tbls15 or len(tbls15[0].rows) != 6 or len(tbls15[0].columns) != 5:
        issues.append("Slide 15 thiếu bảng cluster 6×5")
    # Chart + SHAP: slide 7, 11, 13, 14, 16, 19 mỗi slide ≥ 1 ảnh; tổng ≥ 12 ảnh
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    chart_slides = {6, 10, 12, 13, 15, 18}
    total_pics = 0
    for i, s in enumerate(slides, 1):
        pics = [sh for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        total_pics += len(pics)
        if i - 1 in chart_slides and not pics:
            issues.append(f"Slide {i} thiếu ảnh chart")
        for pic in pics:
            if (pic.left < 0 or pic.top < 0 or
                    pic.left + pic.width > Inches(13.33) or
                    pic.top + pic.height > Inches(7.5)):
                issues.append(f"Slide {i}: ảnh vượt ranh giới slide")
    if total_pics < 12:
        issues.append(f"Tổng ảnh chart = {total_pics}, mong đợi ≥ 12")
    # Ảnh SHAP (2 slide)
    for idx in (12, 16):  # slide 13, 17
        pics = [sh for sh in slides[idx].shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not pics:
            issues.append(f"Slide {idx + 1} thiếu ảnh SHAP")
    # Số liệu chính
    all_text = " ".join(
        sh.text_frame.text for s in slides for sh in s.shapes if sh.has_text_frame).lower()
    for phrase in ["1.193", "4.17", "0.38", "1500", "56%", "6.6%"]:
        if phrase not in all_text:
            issues.append(f"Thiếu số liệu: '{phrase}'")
    if issues:
        print("VERIFICATION FAILED:")
        for i in issues:
            print(f"  - {i}")
        return False
    print(f"VERIFICATION PASSED: {n} slides, bảng đúng kích thước, số liệu đầy đủ!")
    return True


if __name__ == "__main__":
    build()
    verify()
